#!/usr/bin/env python3
"""
Test RAG Agent - Demonstrates document processing without OpenAI API key
This version shows the document processing and chunking capabilities.
"""

import os
import json
from pathlib import Path
from typing import List, Dict, Any
import logging
from datetime import datetime

# Import config
try:  # noqa: SIM105
    from AtendentePro import config
except ModuleNotFoundError:  # pragma: no cover - fallback when running standalone
    import sys
    sys.path.append(str(Path(__file__).resolve().parents[2]))
    from AtendentePro import config  # type: ignore

OPENAI_API_KEY = config.OPENAI_API_KEY

# Document processing imports
import PyPDF2
import docx
from pptx import Presentation
import fitz  # PyMuPDF for better PDF handling

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class TestDocumentProcessor:
    """Handles document processing and text extraction for testing"""
    
    def __init__(self, doc_folder: str = "/Users/arthurvaz/Desktop/Monkai/AtendentePRO/AtendentePro/Template/knowledge_documentos"):
        self.doc_folder = Path(doc_folder)
        self.json_folder = self.doc_folder / "json_format"
        self.documents = {}
        self.chunks = []
        
        # Create json_format directory if it doesn't exist
        self.json_folder.mkdir(exist_ok=True)
        
    def extract_text_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF files"""
        try:
            # Try PyMuPDF first for better text extraction
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
        except Exception as e:
            logger.warning(f"PyMuPDF failed for {file_path}, trying PyPDF2: {e}")
            try:
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in reader.pages:
                        text += page.extract_text()
                    return text
            except Exception as e2:
                logger.error(f"Failed to extract text from {file_path}: {e2}")
                return ""
    
    def extract_text_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint files"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Failed to extract text from {file_path}: {e}")
            return ""
    
    def extract_text_from_docx(self, file_path: Path) -> str:
        """Extract text from Word documents"""
        try:
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Failed to extract text from {file_path}: {e}")
            return ""
    
    def process_documents(self) -> Dict[str, str]:
        """Process all documents in the doc folder"""
        logger.info("Processing documents...")
        
        for file_path in self.doc_folder.iterdir():
            if file_path.is_file():
                file_ext = file_path.suffix.lower()
                file_name = file_path.name
                
                logger.info(f"Processing: {file_name}")
                
                if file_ext == '.pdf':
                    text = self.extract_text_from_pdf(file_path)
                elif file_ext == '.pptx':
                    text = self.extract_text_from_pptx(file_path)
                elif file_ext == '.docx':
                    text = self.extract_text_from_docx(file_path)
                else:
                    logger.warning(f"Unsupported file type: {file_ext}")
                    continue
                
                if text.strip():
                    self.documents[file_name] = text
                    logger.info(f"Successfully processed {file_name} ({len(text)} characters)")
                else:
                    logger.warning(f"No text extracted from {file_name}")
        
        logger.info(f"Processed {len(self.documents)} documents")
        return self.documents
    
    def create_chunks(self, chunk_size: int = 1000, overlap: int = 200) -> List[Dict[str, Any]]:
        """Create overlapping chunks from documents"""
        logger.info("Creating document chunks...")
        
        for doc_name, text in self.documents.items():
            # Simple chunking by character count
            for i in range(0, len(text), chunk_size - overlap):
                chunk = text[i:i + chunk_size]
                if len(chunk.strip()) > 100:  # Only keep meaningful chunks
                    self.chunks.append({
                        'content': chunk,
                        'source': doc_name,
                        'start_pos': i,
                        'end_pos': min(i + chunk_size, len(text))
                    })
        
        logger.info(f"Created {len(self.chunks)} chunks")
        return self.chunks
    
    def analyze_documents(self) -> Dict[str, Any]:
        """Analyze the processed documents and provide insights"""
        analysis = {
            'total_documents': len(self.documents),
            'total_chunks': len(self.chunks),
            'document_types': {},
            'document_sizes': {},
            'sample_chunks': []
        }
        
        # Analyze document types
        for doc_name in self.documents.keys():
            ext = Path(doc_name).suffix.lower()
            analysis['document_types'][ext] = analysis['document_types'].get(ext, 0) + 1
        
        # Analyze document sizes
        for doc_name, text in self.documents.items():
            analysis['document_sizes'][doc_name] = {
                'characters': len(text),
                'words': len(text.split()),
                'lines': len(text.split('\n'))
            }
        
        # Sample chunks
        if self.chunks:
            analysis['sample_chunks'] = [
                {
                    'source': chunk['source'],
                    'content_preview': chunk['content'][:200] + "..." if len(chunk['content']) > 200 else chunk['content'],
                    'length': len(chunk['content'])
                }
                for chunk in self.chunks[:3]  # First 3 chunks
            ]
        
        return analysis
    
    def save_results(self) -> Dict[str, str]:
        """Save processing results to files"""
        # Create a unique identifier based on document names
        doc_names = sorted(list(self.documents.keys()))
        
        # Create a shorter, more readable identifier
        if len(doc_names) == 1:
            # Single document: use its name
            doc_hash = Path(doc_names[0]).stem.replace(" ", "_")[:30]  # Limit to 30 chars
        elif len(doc_names) <= 3:
            # Few documents: combine names
            doc_hash = "_".join([Path(name).stem.replace(" ", "_")[:10] for name in doc_names])
        else:
            # Many documents: use count
            doc_hash = f"{len(doc_names)}_documents"
        
        # Clean up special characters
        doc_hash = "".join(c for c in doc_hash if c.isalnum() or c in ['_', '-'])[:50]
        
        saved_files = {}
        
        try:
            # Save chunks as JSON
            chunks_file = self.json_folder / f"chunks_{doc_hash}.json"
            with open(chunks_file, 'w', encoding='utf-8') as f:
                json.dump(self.chunks, f, ensure_ascii=False, indent=2)
            saved_files['chunks'] = str(chunks_file)
            logger.info(f"Chunks saved to {chunks_file}")
            
            # Save documents as JSON
            documents_file = self.json_folder / f"documents_{doc_hash}.json"
            with open(documents_file, 'w', encoding='utf-8') as f:
                json.dump(self.documents, f, ensure_ascii=False, indent=2)
            saved_files['documents'] = str(documents_file)
            logger.info(f"Documents saved to {documents_file}")
            
            # Save analysis as JSON
            analysis = self.analyze_documents()
            analysis_file = self.json_folder / f"analysis_{doc_hash}.json"
            with open(analysis_file, 'w', encoding='utf-8') as f:
                json.dump(analysis, f, ensure_ascii=False, indent=2)
            saved_files['analysis'] = str(analysis_file)
            logger.info(f"Analysis saved to {analysis_file}")
            
            # Save summary report as text
            report_file = self.json_folder / f"report_{doc_hash}.txt"
            with open(report_file, 'w', encoding='utf-8') as f:
                f.write("RAG Agent Test Results Report\n")
                f.write("=" * 50 + "\n")
                f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
                
                f.write(f"Total Documents Processed: {analysis['total_documents']}\n")
                f.write(f"Total Chunks Created: {analysis['total_chunks']}\n\n")
                
                f.write("Document Types:\n")
                for ext, count in analysis['document_types'].items():
                    f.write(f"  {ext}: {count} files\n")
                f.write("\n")
                
                f.write("Document Sizes:\n")
                for doc_name, stats in analysis['document_sizes'].items():
                    f.write(f"  {doc_name}:\n")
                    f.write(f"    Characters: {stats['characters']:,}\n")
                    f.write(f"    Words: {stats['words']:,}\n")
                    f.write(f"    Lines: {stats['lines']:,}\n\n")
                
                if analysis['sample_chunks']:
                    f.write("Sample Chunks:\n")
                    for i, chunk in enumerate(analysis['sample_chunks'], 1):
                        f.write(f"  Chunk {i} (from {chunk['source']}):\n")
                        f.write(f"    Length: {chunk['length']} characters\n")
                        f.write(f"    Preview: {chunk['content_preview']}\n\n")
            
            saved_files['report'] = str(report_file)
            logger.info(f"Report saved to {report_file}")
            
            # Save latest files (without document hash for easy access)
            latest_chunks = self.json_folder / "chunks_latest.json"
            latest_documents = self.json_folder / "documents_latest.json"
            latest_analysis = self.json_folder / "analysis_latest.json"
            latest_report = self.json_folder / "report_latest.txt"
            
            # Copy to latest files
            import shutil
            shutil.copy2(chunks_file, latest_chunks)
            shutil.copy2(documents_file, latest_documents)
            shutil.copy2(analysis_file, latest_analysis)
            shutil.copy2(report_file, latest_report)
            
            saved_files['latest_chunks'] = str(latest_chunks)
            saved_files['latest_documents'] = str(latest_documents)
            saved_files['latest_analysis'] = str(latest_analysis)
            saved_files['latest_report'] = str(latest_report)
            
            logger.info("Latest files updated")
            
        except Exception as e:
            logger.error(f"Error saving results: {e}")
            raise
        
        return saved_files

def main():
    """Main test function"""
    print("🧪 Testing RAG Agent Document Processing")
    print("=" * 50)
    
    # Check if doc folder exists
    doc_path = Path("/Users/arthurvaz/Desktop/Monkai/AtendentePRO/AtendentePro/Template/knowledge_documentos")
    if not doc_path.exists():
        print("❌ Error: 'knowledge_documentos' folder not found!")
        print("Please ensure you have documents in the knowledge_documentos folder.")
        return
    
    # Check if there are documents in the folder
    doc_files = list(doc_path.glob("*"))
    if not doc_files:
        print("❌ Error: No files found in 'doc' folder!")
        print("Please add some documents (PDF, PPTX, DOCX) to the 'doc' folder.")
        return
    
    print(f"📁 Found {len(doc_files)} files in doc folder:")
    for file in doc_files:
        print(f"   - {file.name}")
    
    print("\n🔄 Processing documents...")
    
    # Initialize processor
    processor = TestDocumentProcessor()
    
    try:
        # Process documents
        documents = processor.process_documents()
        
        if not documents:
            print("❌ No documents were successfully processed!")
            return
        
        print(f"\n✅ Successfully processed {len(documents)} documents!")
        
        # Create chunks
        chunks = processor.create_chunks()
        print(f"📝 Created {len(chunks)} text chunks")
        
        # Analyze results
        analysis = processor.analyze_documents()
        
        print("\n📊 Document Analysis:")
        print(f"   Total Documents: {analysis['total_documents']}")
        print(f"   Total Chunks: {analysis['total_chunks']}")
        
        print("\n📄 Document Types:")
        for ext, count in analysis['document_types'].items():
            print(f"   {ext}: {count} files")
        
        print("\n📏 Document Sizes:")
        for doc_name, stats in analysis['document_sizes'].items():
            print(f"   {doc_name}:")
            print(f"     Characters: {stats['characters']:,}")
            print(f"     Words: {stats['words']:,}")
            print(f"     Lines: {stats['lines']:,}")
        
        if analysis['sample_chunks']:
            print("\n🔍 Sample Chunks:")
            for i, chunk in enumerate(analysis['sample_chunks'], 1):
                print(f"   Chunk {i} (from {chunk['source']}):")
                print(f"     Length: {chunk['length']} characters")
                print(f"     Preview: {chunk['content_preview']}")
                print()
        
        # Save results to files
        print("\n💾 Saving results to files...")
        saved_files = processor.save_results()
        
        print("✅ Document processing test completed successfully!")
        print(f"\n📁 Results saved to: {processor.json_folder}")
        print("Files created:")
        for file_type, file_path in saved_files.items():
            if file_type.startswith('latest_'):
                print(f"   📄 {file_type}: {Path(file_path).name}")
        
        # Check if API key is available
        if OPENAI_API_KEY and OPENAI_API_KEY.strip():
            print(f"\n🔑 OpenAI API Key encontrada no config.py!")
            print("💡 Você pode usar o RAG Agent completo com IA:")
            print("   Run: python rag_agent.py")
            
            # Offer to run the full RAG agent
            try:
                response = input("\n🤖 Deseja executar o RAG Agent completo agora? (s/n): ").strip().lower()
                if response in ['s', 'sim', 'y', 'yes']:
                    print("\n🚀 Iniciando RAG Agent completo...")
                    # Import and run the full RAG agent
                    from rag_agent import main as rag_main
                    import asyncio
                    asyncio.run(rag_main())
            except KeyboardInterrupt:
                print("\n❌ Operação cancelada pelo usuário.")
            except Exception as e:
                print(f"❌ Erro ao executar RAG Agent: {e}")
        else:
            print("\n❌ OpenAI API Key não encontrada no config.py")
            print("💡 Para usar o RAG Agent completo:")
            print("   1. Configure sua API key no arquivo config.py")
            print("   2. Run: python rag_agent.py")
        
    except Exception as e:
        print(f"❌ Error during processing: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()
