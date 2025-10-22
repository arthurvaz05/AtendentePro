#!/usr/bin/env python3
"""
RAG Agent using OpenAI Swarm Framework
This agent processes documents from the doc folder and provides intelligent answers using RAG.
"""

import json
import asyncio
import io
import time
from pathlib import Path
from typing import List, Dict, Any, Optional
import logging
import os
# Document processing imports
import PyPDF2
import docx
from pptx import Presentation
import fitz  # PyMuPDF for better PDF handling
from openai import AzureOpenAI

# Vector database and embeddings
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
import pickle

from AtendentePro.utils.openai_client import get_async_client, get_provider
import config

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Handles document processing, OCR, and text extraction"""

    MIN_OCR_DIM = 200
    MAX_OCR_DIM = 16000
    OCR_RETRY_DELAY_SEC = 10
    SUPPORTED_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff"}

    def __init__(
        self,
        doc_folder: str = "/home/davi/Desktop/AtendentePro/AtendentePro/Template/White_Martins/knowledge_documentos",
        ocr_enabled: Optional[bool] = True,
    ):
        self.doc_folder = Path(doc_folder)
        self.embedding_folder = self.doc_folder / "embedding"
        self.documents = {}
        self.chunks = []

        self.azure_endpoint = getattr(config, "AZURE_AI_VISION_ENDPOINT", None)
        self.azure_key = getattr(config, "AZURE_AI_VISION_KEY", None)
        self.ocr_enabled = self._resolve_ocr_enabled(ocr_enabled)
        self._ocr_client = None

        # Create embedding directory if it doesn't exist
        self.embedding_folder.mkdir(exist_ok=True)

    def _resolve_ocr_enabled(self, explicit_flag: Optional[bool]) -> bool:
        if explicit_flag is not None:
            if explicit_flag and not (self.azure_endpoint and self.azure_key):
                logger.warning("OCR requested but Azure Vision credentials are missing; disabling OCR")
                return False
            return explicit_flag
        return bool(self.azure_endpoint and self.azure_key)

    def _ensure_ocr_client(self):
        if self._ocr_client is not None:
            return self._ocr_client
        if not self.ocr_enabled:
            return None
        if not (self.azure_endpoint and self.azure_key):
            logger.warning("Azure Vision credentials not configured; disabling OCR")
            self.ocr_enabled = False
            return None
        try:
            from azure.ai.vision.imageanalysis import ImageAnalysisClient
            from azure.core.credentials import AzureKeyCredential
        except ImportError as exc:
            logger.error("Azure AI Vision dependencies missing: %s", exc)
            self.ocr_enabled = False
            return None

        self._ocr_client = ImageAnalysisClient(
            endpoint=self.azure_endpoint,
            credential=AzureKeyCredential(self.azure_key),
        )
        return self._ocr_client

    @staticmethod
    def _ensure_allowed_dimensions(img):
        from PIL import Image

        w, h = img.size
        if max(w, h) > DocumentProcessor.MAX_OCR_DIM:
            scale = DocumentProcessor.MAX_OCR_DIM / float(max(w, h))
            new_w = max(1, int(round(w * scale)))
            new_h = max(1, int(round(h * scale)))
            img = img.resize((new_w, new_h), Image.LANCZOS)
        return img

    @staticmethod
    def _encode_image_for_api(img, image_format: str) -> bytes:
        fmt_map = {
            "jpg": "JPEG",
            "jpeg": "JPEG",
            "png": "PNG",
            "bmp": "BMP",
            "tif": "TIFF",
            "tiff": "TIFF",
        }
        pil_fmt = fmt_map.get((image_format or "").lower(), "PNG")
        buffer = io.BytesIO()
        img.save(buffer, format=pil_fmt)
        return buffer.getvalue()

    def _perform_ocr(self, image_bytes: bytes, image_format: str) -> Optional[str]:
        if not self.ocr_enabled:
            return None

        client = self._ensure_ocr_client()
        if client is None:
            return None

        from PIL import Image
        from azure.ai.vision.imageanalysis.models import VisualFeatures
        from azure.core.exceptions import HttpResponseError

        for attempt in range(3):
            try:
                pil_img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                w, h = pil_img.size
                if min(w, h) < self.MIN_OCR_DIM:
                    logger.debug("Skipping OCR: image too small (%sx%s)", w, h)
                    return None

                pil_img = self._ensure_allowed_dimensions(pil_img)
                formatted_bytes = self._encode_image_for_api(pil_img, image_format)

                result = client.analyze(
                    image_data=formatted_bytes,
                    visual_features=[VisualFeatures.READ],
                    language="en",
                    gender_neutral_caption=True,
                )

                lines = []
                if result.read:
                    for block in result.read.blocks:
                        for text_line in block.lines:
                            lines.append(text_line.text)

                ocr_text = "\n".join(lines).strip()
                if ocr_text:
                    return ocr_text
                logger.debug("OCR returned no text for image")
                return None

            except HttpResponseError as exc:
                if exc.status_code == 429 and attempt < 2:
                    logger.warning("Azure Vision rate limit hit; retrying OCR for image")
                    time.sleep(self.OCR_RETRY_DELAY_SEC)
                    continue
                if exc.status_code == 401:
                    logger.error("Azure Vision authentication failed; disabling OCR")
                    self.ocr_enabled = False
                    return None
                logger.error("Azure Vision HTTP error (%s): %s", exc.status_code, exc.message)
                return None
            except Exception as exc:
                logger.error("Unexpected OCR failure: %s", exc)
                return None

        logger.error("OCR failed after maximum retries")
        return None

    def _extract_ocr_from_page(self, pdf_document: fitz.Document, page_index: int) -> List[str]:
        page = pdf_document.load_page(page_index)
        image_list = page.get_images(full=True)
        ocr_snippets: List[str] = []

        for image_index, image in enumerate(image_list):
            try:
                base_image = pdf_document.extract_image(image[0])
                ocr_text = self._perform_ocr(base_image.get("image"), base_image.get("ext", ""))
                if ocr_text:
                    ocr_snippets.append(f"[Page {page_index + 1} Image {image_index + 1}] {ocr_text}")
            except Exception as exc:
                logger.error(
                    "Failed to OCR image %s on page %s (%s): %s",
                    image_index,
                    page_index,
                    pdf_document.name if hasattr(pdf_document, "name") else "<memory>",
                    exc,
                )
        return ocr_snippets

    def extract_text_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF files including OCR content"""
        try:
            text_fragments: List[str] = []
            with fitz.open(file_path) as doc:
                for page_index in range(doc.page_count):
                    page = doc.load_page(page_index)
                    text_fragments.append(page.get_text())
                    if self.ocr_enabled:
                        ocr_fragments = self._extract_ocr_from_page(doc, page_index)
                        text_fragments.extend(ocr_fragments)
            return "\n\n".join(fragment for fragment in text_fragments if fragment)
        except Exception as e:
            logger.warning(f"PyMuPDF failed for {file_path}, trying PyPDF2: {e}")
            try:
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in reader.pages:
                        extracted = page.extract_text()
                        if extracted:
                            text += extracted
                    return text
            except Exception as e2:
                logger.error(f"Failed to extract text from {file_path}: {e2}")
                return ""

    def extract_text_from_image(self, file_path: Path) -> str:
        if not self.ocr_enabled:
            logger.warning("Skipping image %s because OCR is disabled", file_path.name)
            return ""
        try:
            with open(file_path, "rb") as image_file:
                image_bytes = image_file.read()
            text = self._perform_ocr(image_bytes, file_path.suffix.lstrip('.'))
            return text or ""
        except Exception as exc:
            logger.error("Failed to OCR image %s: %s", file_path.name, exc)
            return ""
    
    def extract_text_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint files"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Failed to extract text from {file_path}: {e}")
            return ""
    
    def extract_text_from_docx(self, file_path: Path) -> str:
        """Extract text from Word documents"""
        try:
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Failed to extract text from {file_path}: {e}")
            return ""
    
    def process_documents(self) -> Dict[str, str]:
        """Process all documents in the doc folder"""
        logger.info("Processing documents...")
        
        for file_path in self.doc_folder.iterdir():
            if file_path.is_file():
                file_ext = file_path.suffix.lower()
                file_name = file_path.name
                
                logger.info(f"Processing: {file_name}")
                
                if file_ext == '.pdf':
                    text = self.extract_text_from_pdf(file_path)
                elif file_ext == '.pptx':
                    text = self.extract_text_from_pptx(file_path)
                elif file_ext == '.docx':
                    text = self.extract_text_from_docx(file_path)
                elif file_ext in self.SUPPORTED_IMAGE_EXTENSIONS:
                    text = self.extract_text_from_image(file_path)
                else:
                    logger.warning(f"Unsupported file type: {file_ext}")
                    continue
                
                text = text or ""

                if text.strip():
                    self.documents[file_name] = text
                    logger.info(f"Successfully processed {file_name} ({len(text)} characters)")
                else:
                    logger.warning(f"No text extracted from {file_name}")
        
        logger.info(f"Processed {len(self.documents)} documents")
        return self.documents
    
    def create_chunks(self, chunk_size: int = 1000, overlap: int = 200) -> List[Dict[str, Any]]:
        """Create overlapping chunks from documents"""
        logger.info("Creating document chunks...")
        
        for doc_name, text in self.documents.items():
            # Simple chunking by character count
            for i in range(0, len(text), chunk_size - overlap):
                chunk = text[i:i + chunk_size]
                if len(chunk.strip()) > 100:  # Only keep meaningful chunks
                    self.chunks.append({
                        'content': chunk,
                        'source': doc_name,
                        'start_pos': i,
                        'end_pos': min(i + chunk_size, len(text))
                    })
        
        logger.info(f"Created {len(self.chunks)} chunks")
        return self.chunks

class RAGAgent:
    """Main RAG agent using OpenAI Swarm agentframework"""
    
    def __init__(self, client: Optional[object] = None):
        # Import config here to avoid circular imports
        from config import DEFAULT_MODEL

        self.client = AzureOpenAI(
            api_key=os.getenv("OPENAI_API_KEY_BRASILSOUTH"),
            azure_endpoint=os.getenv("OPENAI_AZURE_ENDPOINT_BRASILSOUTH"),
            api_version=os.getenv("OPENAI_API_VERSION"),
        )
        self.provider = get_provider()
        self.default_model = DEFAULT_MODEL
        
        # Initialize document processor
        self.doc_processor = DocumentProcessor()
        self.embeddings = {}
        self.chunk_embeddings = []
        
    async def process_and_embed_documents(self):
        """Process documents and create embeddings"""
        logger.info("Processing and embedding documents...")
        
        if self.doc_processor.ocr_enabled:
            logger.info(
                "OCR enabled: image content will be processed via Azure Vision endpoint %s",
                getattr(self.doc_processor, "azure_endpoint", "<not configured>"),
            )
        else:
            logger.info("OCR disabled: only textual document content will be embedded")

        # Process documents
        self.doc_processor.process_documents()
        self.doc_processor.create_chunks()
        
        # Create embeddings for chunks
        for i, chunk in enumerate(self.doc_processor.chunks):
            try:
                response = self.client.embeddings.create(
                    model="text-embedding-3-large",
                    input=chunk['content']
                )
                embedding = response.data[0].embedding
                self.chunk_embeddings.append({
                    'chunk': chunk,
                    'embedding': embedding,
                    'index': i
                })
                logger.info(f"Embedded chunk {i+1}/{len(self.doc_processor.chunks)}")
            except Exception as e:
                logger.error(f"Failed to embed chunk {i}: {e}")
        
        logger.info(f"Successfully embedded {len(self.chunk_embeddings)} chunks")
    
    async def find_relevant_chunks(self, query: str, top_k: int = 5) -> List[Dict[str, Any]]:
        """Find most relevant chunks for a given query"""
        try:
            # Get query embedding
            response = self.client.embeddings.create(
                model="text-embedding-3-large",
                input=query
            )
            query_embedding = response.data[0].embedding
            
            # Calculate similarities
            similarities = []
            for chunk_data in self.chunk_embeddings:
                similarity = cosine_similarity(
                    [query_embedding], 
                    [chunk_data['embedding']]
                )[0][0]
                similarities.append((similarity, chunk_data))
            
            # Sort by similarity and return top_k
            similarities.sort(key=lambda x: x[0], reverse=True)
            return [chunk_data for _, chunk_data in similarities[:top_k]]
            
        except Exception as e:
            logger.error(f"Error finding relevant chunks: {e}")
            return []
    
    async def answer_question(self, question: str) -> Dict[str, Any]:
        """Answer a question using RAG"""
        logger.info(f"Processing question: {question}")
        
        # Find relevant chunks
        relevant_chunks = await self.find_relevant_chunks(question, top_k=5)
        
        if not relevant_chunks:
            return {
                'answer': 'I could not find relevant information in the documents to answer your question.',
                'sources': [],
                'confidence': 0.0
            }
        
        # Prepare context for the AI
        context = "\n\n".join([
            f"Document: {chunk['chunk']['source']}\n"
            f"Content: {chunk['chunk']['content']}"
            for chunk in relevant_chunks
        ])
        
        # Create prompt for the AI
        prompt = f"""
        Based on the following document excerpts, please answer the user's question.
        If the information is not available in the provided context, say so.
        Provide specific references to the source documents when possible.
        
        Question: {question}
        
        Document Context:
        {context}
        
        Please provide a comprehensive answer based on the available information.
        """
        
        try:
            response = await self.client.chat.completions.create(
                model=self.default_model,
                messages=[
                    {"role": "system", "content": """Você é um agente RAG (Geração Aumentada por Recuperação) especializado em equipamentos e materiais de soldagem.
                                                        Seu papel é responder perguntas com base no contexto documental fornecido.
                                                        Sempre forneça respostas precisas e úteis, citando os documentos de origem sempre que possível.
                                                        Se a informação não estiver disponível no contexto, declare isso claramente."""},
                    {"role": "user", "content": prompt},
                ],
            )
            # Extract the answer from the response
            answer = response.choices[0].message.content
            
            return {
                'answer': answer,
                'sources': [chunk['chunk']['source'] for chunk in relevant_chunks],
                'confidence': 0.8,  # Placeholder confidence score
                'context_used': context[:500] + "..." if len(context) > 500 else context
            }
            
        except Exception as e:
            logger.error(f"Error generating answer: {e}")
            return {
                'answer': f'I encountered an error while processing your question: {str(e)}',
                'sources': [],
                'confidence': 0.0
            }
    
    def save_embeddings(self, filename: str = "embeddings.pkl"):
        """Save embeddings to disk for reuse"""
        try:
            filepath = self.doc_processor.embedding_folder / filename
            with open(filepath, 'wb') as f:
                pickle.dump(self.chunk_embeddings, f)
            logger.info(f"Embeddings saved to {filepath}")
        except Exception as e:
            logger.error(f"Failed to save embeddings: {e}")
    
    def load_embeddings(self, filename: str = "embeddings.pkl"):
        """Load embeddings from disk"""
        try:
            filepath = self.doc_processor.embedding_folder / filename
            with open(filepath, 'rb') as f:
                self.chunk_embeddings = pickle.load(f)
            logger.info(f"Embeddings loaded from {filepath}")
        except Exception as e:
            logger.error(f"Failed to load embeddings: {e}")

async def main():
    """Main function to demonstrate the RAG agent"""
    
    provider = get_provider()
    print(f"Using OpenAI provider: {provider}")
    
    # Initialize the RAG agent
    try:
        agent = RAGAgent()
    except RuntimeError as exc:
        print(f"Please configure your OpenAI credentials: {exc}")
        return
    
    # Process documents (this will take some time)
    print("Processing documents and creating embeddings...")
    await agent.process_and_embed_documents()
    
    # Save embeddings for future use
    agent.save_embeddings()
    
    # Interactive question answering
    print("\nRAG Agent is ready! Ask questions about the documents.")
    print("Type 'quit' to exit.\n")
    
    while True:
        try:
            question = input("Your question: ").strip()
            if question.lower() in ['quit', 'exit', 'q']:
                break
            
            if not question:
                continue
            
            print("\nProcessing your question...")
            result = await agent.answer_question(question)
            
            print(f"\nAnswer: {result['answer']}")
            print(f"Sources: {', '.join(result['sources'])}")
            print(f"Confidence: {result['confidence']:.2f}")
            print("-" * 80)
            
        except KeyboardInterrupt:
            break
        except Exception as e:
            print(f"Error: {e}")
    
    print("Goodbye!")

if __name__ == "__main__":
    # Install required packages if not available
    try:
        import PyPDF2
        import docx
        from pptx import Presentation
        import fitz
        import sklearn
    except ImportError as e:
        print(f"Missing required package: {e}")
        print("Please install required packages:")
        print("pip install PyPDF2 python-docx python-pptx PyMuPDF scikit-learn")
        exit(1)
    
    # Run the main function
    asyncio.run(main())
